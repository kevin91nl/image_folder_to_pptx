'''
This is the main file that runs the program.
'''
import tkinter as tk
from tkinter import filedialog
import os
from pptx import Presentation
from pptx.util import Inches
def select_folder():
    """
    Opens a file dialog to select a folder containing photos.
    """
    folder_path = filedialog.askdirectory()
    if folder_path:
        create_presentation(folder_path)
    else:
        # Handle the case when the user cancels the folder selection
        print("No folder selected. Exiting the program.")
def create_presentation(folder_path):
    """
    Creates a PowerPoint presentation and inserts photos from the selected folder.
    Allows the user to specify the save location and filename.
    """
    prs = Presentation()
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            if file.endswith(('.jpg', '.jpeg', '.png')):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                left = Inches(1)
                top = Inches(1)
                width = Inches(8)
                height = Inches(6)
                slide.shapes.add_picture(os.path.join(root, file), left, top, width, height)
    # Prompt the user to select the save location and filename
    save_path = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint Presentation", "*.pptx")])
    if save_path:
        prs.save(save_path)
def main():
    """
    Main function that initializes the program.
    """
    root = tk.Tk()
    root.withdraw()
    select_folder()
if __name__ == "__main__":
    main()